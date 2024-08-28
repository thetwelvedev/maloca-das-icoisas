from pptx import Presentation

def ler_pptx(caminho_arquivo):
    # Carrega a apresentação
    prs = Presentation(caminho_arquivo)
    
    # Itera sobre os slides
    for i, slide in enumerate(prs.slides):
        print(f"Slide {i + 1}")
        
        # Itera sobre os shapes (elementos) no slide
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
                
if __name__ == "__main__":
    # Caminho para o arquivo .pptx
    caminho_arquivo = "exemplo.pptx"
    ler_pptx(caminho_arquivo)
